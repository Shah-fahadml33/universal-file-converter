from pathlib import Path

import fitz
import streamlit as st
from PIL import Image
from docx import Document
from pptx import Presentation


def preview_pdf(pdf_path: str):
    """
    Display the first page of a PDF.
    """

    document = fitz.open(pdf_path)

    if len(document) == 0:
        st.warning("Empty PDF.")
        return

    page = document.load_page(0)

    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))

    image_path = Path(pdf_path).with_suffix(".preview.png")

    pix.save(image_path)

    st.image(str(image_path), caption="PDF Preview")

    document.close()

    image_path.unlink(missing_ok=True)


def preview_docx(docx_path: str):
    """
    Display first paragraphs of a DOCX.
    """

    document = Document(docx_path)

    st.subheader("Document Preview")

    for paragraph in document.paragraphs[:10]:
        if paragraph.text.strip():
            st.write(paragraph.text)


def preview_image(image_path: str):
    """
    Display an image.
    """

    image = Image.open(image_path)

    st.image(image)

    st.write(f"Width : {image.width}")
    st.write(f"Height : {image.height}")
    st.write(f"Format : {image.format}")
    st.write(f"Mode : {image.mode}")


def preview_pptx(pptx_path: str):
    """
    Display PowerPoint information.
    """

    prs = Presentation(pptx_path)

    st.subheader("Presentation Preview")

    st.write(f"Slides : {len(prs.slides)}")

    for index, slide in enumerate(prs.slides[:3], start=1):

        st.markdown(f"### Slide {index}")

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                if shape.text.strip():
                    st.write(shape.text)


def preview(file_path: str):
    """
    Automatically preview supported files.
    """

    extension = Path(file_path).suffix.lower()

    if extension == ".pdf":
        preview_pdf(file_path)

    elif extension == ".docx":
        preview_docx(file_path)

    elif extension in [".png", ".jpg", ".jpeg", ".bmp", ".webp"]:
        preview_image(file_path)

    elif extension == ".pptx":
        preview_pptx(file_path)

    else:
        st.info("Preview not available.")