from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image


def create_blank_presentation(output_path: str):
    """
    Create an empty PowerPoint presentation.
    """

    prs = Presentation()
    prs.save(output_path)

    return output_path


def images_to_pptx(image_paths, output_path):
    """
    Convert multiple images into a PowerPoint presentation.
    """

    prs = Presentation()

    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    for image_path in image_paths:

        slide = prs.slides.add_slide(blank_layout)

        slide.shapes.add_picture(
            image_path,
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(output_path)

    return output_path


def extract_slide_text(pptx_path: str):
    """
    Extract all text from a PowerPoint presentation.
    """

    prs = Presentation(pptx_path)

    text = []

    for slide in prs.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):
                if shape.text.strip():
                    text.append(shape.text)

    return "\n".join(text)


def get_presentation_info(pptx_path: str):
    """
    Return basic presentation information.
    """

    prs = Presentation(pptx_path)

    return {
        "slides": len(prs.slides),
        "width": prs.slide_width,
        "height": prs.slide_height,
    }


def validate_pptx(pptx_path: str):
    """
    Validate a PowerPoint file.
    """

    try:
        Presentation(pptx_path)
        return True
    except Exception:
        return False


def image_to_single_slide(
    image_path: str,
    output_path: str
):
    """
    Create a PowerPoint with one image.
    """

    prs = Presentation()

    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)

    slide.shapes.add_picture(
        image_path,
        0,
        0,
        width=prs.slide_width,
        height=prs.slide_height,
    )

    prs.save(output_path)

    return output_path