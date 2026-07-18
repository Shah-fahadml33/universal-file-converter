from pathlib import Path

import fitz  # PyMuPDF
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Inches


def pdf_to_docx(pdf_path: str, output_path: str):
    """
    Convert PDF to DOCX.
    """

    cv = Converter(pdf_path)

    try:
        cv.convert(output_path)
    finally:
        cv.close()

    return output_path


def pdf_to_images(pdf_path: str, output_folder: str):
    """
    Convert each page of a PDF into a PNG image.
    Returns a list of image paths.
    """

    Path(output_folder).mkdir(parents=True, exist_ok=True)

    document = fitz.open(pdf_path)

    image_paths = []

    for page_number in range(len(document)):

        page = document.load_page(page_number)

        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))

        image_path = Path(output_folder) / f"page_{page_number + 1}.png"

        pix.save(str(image_path))

        image_paths.append(str(image_path))

    document.close()

    return image_paths


def pdf_to_pptx(pdf_path: str, output_path: str):
    """
    Convert PDF pages into PowerPoint slides.
    """

    image_folder = Path(output_path).parent / "temp_images"

    images = pdf_to_images(pdf_path, image_folder)

    prs = Presentation()

    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    for image in images:

        slide = prs.slides.add_slide(blank_layout)

        slide.shapes.add_picture(
            image,
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(output_path)

    return output_path


def merge_pdfs(pdf_files, output_path):
    """
    Merge multiple PDFs into one.
    """

    merged = fitz.open()

    for pdf in pdf_files:
        doc = fitz.open(pdf)
        merged.insert_pdf(doc)
        doc.close()

    merged.save(output_path)
    merged.close()

    return output_path


def split_pdf(pdf_path, output_folder):
    """
    Split a PDF into one PDF per page.
    """

    Path(output_folder).mkdir(parents=True, exist_ok=True)

    document = fitz.open(pdf_path)

    output_files = []

    for page_num in range(len(document)):

        single = fitz.open()

        single.insert_pdf(
            document,
            from_page=page_num,
            to_page=page_num,
        )

        filename = Path(output_folder) / f"page_{page_num + 1}.pdf"

        single.save(filename)

        single.close()

        output_files.append(str(filename))

    document.close()

    return output_files


def rotate_pdf(pdf_path, output_path, rotation=90):
    """
    Rotate every page in a PDF.
    """

    document = fitz.open(pdf_path)

    for page in document:
        page.set_rotation(rotation)

    document.save(output_path)

    document.close()

    return output_path


def compress_pdf(pdf_path, output_path):
    """
    Save a PDF with garbage collection and compression.
    """

    document = fitz.open(pdf_path)

    document.save(
        output_path,
        garbage=4,
        deflate=True,
        clean=True,
    )

    document.close()

    return output_path